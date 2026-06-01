from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "IvyEuro_Education_网站改造分析报告.pptx"

BLUE = RGBColor(27, 109, 255)
NAVY = RGBColor(11, 18, 32)
GREEN = RGBColor(34, 163, 92)
TEAL = RGBColor(25, 196, 176)
AMBER = RGBColor(245, 158, 11)
RED = RGBColor(239, 68, 68)
MUTED = RGBColor(90, 107, 127)
TEXT = RGBColor(16, 32, 51)
LIGHT = RGBColor(244, 247, 251)
WHITE = RGBColor(255, 255, 255)


def set_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header(slide, title, subtitle=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.92))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    t = slide.shapes.add_textbox(Inches(0.55), Inches(0.18), Inches(10.6), Inches(0.42))
    tf = t.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = WHITE

    if subtitle:
        st = slide.shapes.add_textbox(Inches(9.2), Inches(0.2), Inches(3.6), Inches(0.34))
        st.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        r2 = st.text_frame.paragraphs[0].add_run()
        r2.text = subtitle
        r2.font.name = "Microsoft YaHei"
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = RGBColor(204, 214, 224)


def add_title_block(slide, title, subtitle, x=0.65, y=1.15, w=11.8):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(1.4))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = TEXT

    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = "Microsoft YaHei"
        r2.font.size = Pt(12.5)
        r2.font.color.rgb = MUTED
        p2.space_before = Pt(10)


def add_card(slide, x, y, w, h, title, body, fill=LIGHT, accent=BLUE, title_size=15, body_size=11):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RGBColor(220, 227, 235)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.12), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(x + 0.25), Inches(y + 0.18), Inches(w - 0.4), Inches(h - 0.3))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(title_size)
    r.font.bold = True
    r.font.color.rgb = TEXT
    p.space_after = Pt(4)

    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = body
    r2.font.name = "Microsoft YaHei"
    r2.font.size = Pt(body_size)
    r2.font.color.rgb = MUTED
    p2.space_before = Pt(4)


def add_box(slide, x, y, w, h, text, fill=WHITE, line=BLUE, font_size=14, bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = 1
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = TEXT
    return shp


def add_arrow(slide, x1, y1, x2, y2, color=RGBColor(120, 130, 145)):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(1.6)
    conn.line.end_arrowhead = True
    return conn


def add_bullet_list(slide, x, y, w, h, title, bullets, accent=BLUE):
    add_card(slide, x, y, w, h, title, "", fill=WHITE, accent=accent)
    tx = slide.shapes.add_textbox(Inches(x + 0.25), Inches(y + 0.62), Inches(w - 0.45), Inches(h - 0.75))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.bullet = True
        p.space_after = Pt(7)
        r = p.add_run()
        r.text = b
        r.font.name = "Microsoft YaHei"
        r.font.size = Pt(12)
        r.font.color.rgb = TEXT


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Slide 1 cover
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, NAVY)

bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = NAVY
bg.line.fill.background()

accent = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(0.72), Inches(3.1), Inches(0.48))
accent.fill.solid()
accent.fill.fore_color.rgb = BLUE
accent.line.fill.background()
accent.text_frame.text = "官网改造 / 内容迁移 / PPT汇报"
accent.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
for r in accent.text_frame.paragraphs[0].runs:
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = WHITE

tb = slide.shapes.add_textbox(Inches(0.75), Inches(1.35), Inches(8.8), Inches(2.0))
tf = tb.text_frame
p = tf.paragraphs[0]
r = p.add_run()
r.text = "IvyEuro Education\n网站改造分析报告"
r.font.name = "Microsoft YaHei"
r.font.size = Pt(30)
r.font.bold = True
r.font.color.rgb = WHITE

p2 = tf.add_paragraph()
r2 = p2.add_run()
r2.text = "从老站内容迁移到多页面招生官网，完成首页、课程、新闻、联系与移动端体验重构。"
r2.font.name = "Microsoft YaHei"
r2.font.size = Pt(14)
r2.font.color.rgb = RGBColor(204, 214, 224)
p2.space_before = Pt(12)

meta = slide.shapes.add_textbox(Inches(0.75), Inches(4.0), Inches(5.6), Inches(1.2))
mtf = meta.text_frame
for i, line in enumerate([
    "项目地址：guoyunyi0417-hue / ivy-euro-education",
    f"报告日期：{date.today().isoformat()}",
    "重点：内容迁移、视觉升级、联系方式统一",
]):
    p = mtf.paragraphs[0] if i == 0 else mtf.add_paragraph()
    r = p.add_run()
    r.text = line
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(12.5)
    r.font.color.rgb = RGBColor(225, 232, 241)
    p.space_after = Pt(6)

badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.25), Inches(1.0), Inches(2.25), Inches(2.25))
badge.fill.solid()
badge.fill.fore_color.rgb = RGBColor(18, 31, 52)
badge.line.color.rgb = BLUE
badge.line.width = Pt(1.5)
bt = badge.text_frame
bt.vertical_anchor = 1
bp = bt.paragraphs[0]
bp.alignment = PP_ALIGN.CENTER
rr = bp.add_run()
rr.text = "Barcelonа\nOnline School"
rr.font.name = "Microsoft YaHei"
rr.font.size = Pt(16)
rr.font.bold = True
rr.font.color.rgb = WHITE

add_card(slide, 8.55, 4.05, 3.8, 1.65, "改造关键词", "多页面结构 · 招生官网 · 课程表 · 新闻动态 · 微信 / WhatsApp 联系", fill=RGBColor(18, 31, 52), accent=TEAL, title_size=14, body_size=12)

# Slide 2 project overview
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_header(slide, "项目概况", "页面化结构 + 信息重组")
add_title_block(slide, "为什么要改造", "旧站内容有价值，但结构过散、视觉较旧、转化入口不统一。")
add_bullet_list(slide, 0.65, 2.05, 5.95, 4.95, "原站主要问题", [
    "信息分散，首页与栏目缺少统一导航",
    "老师、新闻、好评页多为占位式内容",
    "暑期课表适合短期展示，不适合长期压在首页",
    "联系入口偏弱，缺少微信 / WhatsApp 的明确动作",
])
add_bullet_list(slide, 6.7, 2.05, 5.95, 4.95, "改造目标", [
    "把老内容迁移成可持续更新的招生官网",
    "首页突出课程、咨询、课程表和转化入口",
    "让新闻和课程形成独立栏目，便于后续维护",
    "补齐微信、WhatsApp、二维码和表单等联系路径",
], accent=GREEN)

# Slide 3 site map
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_header(slide, "网站结构图", "从单页展示升级为多页面官网")
add_title_block(slide, "页面结构关系", "首页是转化入口，栏目页负责内容承接，新闻详情用于深度阅读。")

home = add_box(slide, 5.2, 1.95, 2.95, 0.72, "首页", fill=RGBColor(235, 243, 255))
labels = [
    ("服务项目", 0.6, 3.0),
    ("关于我们", 2.45, 3.0),
    ("老师简介", 4.3, 3.0),
    ("课程表", 6.15, 3.0),
    ("家长好评", 8.0, 3.0),
    ("新闻动态", 9.85, 3.0),
    ("联系我们", 11.7, 3.0),
]
targets = []
for text, x, y in labels:
    targets.append(add_box(slide, x, y, 1.45, 0.65, text, fill=WHITE, line=BLUE, font_size=12))
    add_arrow(slide, 6.6, 2.67, x + 0.72, 3.0, color=RGBColor(145, 155, 170))

news_detail = add_box(slide, 10.0, 4.25, 1.95, 0.62, "新闻详情", fill=RGBColor(238, 250, 245), line=GREEN, font_size=12)
add_arrow(slide, 10.6, 3.65, 10.95, 4.25, color=GREEN)

note = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(5.35), Inches(11.9), Inches(1.0))
note.fill.solid()
note.fill.fore_color.rgb = LIGHT
note.line.color.rgb = RGBColor(220, 227, 235)
nt = note.text_frame
np = nt.paragraphs[0]
nr = np.add_run()
nr.text = "结构原则：长期稳定信息放在首页和栏目页，季节性内容放在课程表 / 公告位，新闻栏目支持后续持续更新。"
nr.font.name = "Microsoft YaHei"
nr.font.size = Pt(13)
nr.font.color.rgb = TEXT
np.alignment = PP_ALIGN.CENTER

# Slide 4 content assessment matrix
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_header(slide, "老资料内容评估", "保留 / 替换 / 补强 / 删除四类")
add_title_block(slide, "内容处置图", "先保留业务信息，再替换模板句，最后补强真实案例。")

quadrants = [
    ("保留", "机构名称、地址、电话、营业时间\n课程大类、微信 / WhatsApp、二维码", 0.65, 2.0, GREEN),
    ("替换", "老师占位卡、家长好评模板句\n新闻/文章标题、长期首页暑期大块", 6.85, 2.0, AMBER),
    ("补强", "真实老师介绍、真实家长评价\n学习成果、活动照片、真实通知", 0.65, 4.55, BLUE),
    ("删除 / 临时", "过期的暑期公告、首页大块课表\n没有实际内容支撑的草稿句", 6.85, 4.55, RED),
]
for t, b, x, y, c in quadrants:
    add_card(slide, x, y, 5.75, 1.85, t, b, fill=WHITE, accent=c, title_size=16, body_size=12)

bottom = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(6.55), Inches(11.9), Inches(0.55))
bottom.fill.solid()
bottom.fill.fore_color.rgb = RGBColor(242, 246, 250)
bottom.line.color.rgb = RGBColor(220, 227, 235)
tf = bottom.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "判断标准：有事实依据的保留，有营销价值的润色，没有实际内容支撑的占位句要替换。"
r.font.name = "Microsoft YaHei"
r.font.size = Pt(12.5)
r.font.color.rgb = TEXT

# Slide 5 homepage screenshot and advantages
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_header(slide, "首页优势图示", "视觉升级 + 转化入口 + 季节公告")
add_title_block(slide, "新首页的优势", "首页不再只是信息堆砌，而是能承接咨询和转化。")

slide.shapes.add_picture(str(ROOT / "report-assets" / "home-hero.png"), Inches(0.6), Inches(2.0), width=Inches(7.45))
add_card(slide, 8.25, 1.95, 4.35, 1.25, "优势 1：首屏更清晰", "把学校定位、课程亮点、免费测试和联系入口集中到首页首屏。", fill=LIGHT, accent=BLUE, title_size=14, body_size=11)
add_card(slide, 8.25, 3.35, 4.35, 1.25, "优势 2：视觉更像官网", "通过大图、卡片和统一配色，形成更强的机构品牌感。", fill=LIGHT, accent=GREEN, title_size=14, body_size=11)
add_card(slide, 8.25, 4.75, 4.35, 1.25, "优势 3：临时内容可自动消失", "暑期班只做季节公告，到期后自动隐藏，不污染首页长期信息。", fill=LIGHT, accent=AMBER, title_size=14, body_size=11)
add_card(slide, 8.25, 6.15, 4.35, 0.95, "优势 4：移动端体验更统一", "导航、联系方式、新闻栏目都支持手机浏览。", fill=LIGHT, accent=TEAL, title_size=13, body_size=10)

# Slide 6 contact channels
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_header(slide, "联系通道图", "微信优先，WhatsApp 备用")
add_title_block(slide, "家长如何联系学校", "把核心联系方式做成明确入口，减少家长的操作成本。")

add_box(slide, 0.75, 2.35, 2.15, 0.8, "微信昵称\nIvyEuro学校刘老师", fill=RGBColor(241, 250, 241), line=GREEN, font_size=12)
add_box(slide, 3.2, 2.35, 2.0, 0.8, "微信号\nsophialiu0826", fill=RGBColor(241, 250, 241), line=GREEN, font_size=12)
add_box(slide, 5.5, 2.35, 1.6, 0.8, "二维码", fill=RGBColor(241, 250, 241), line=GREEN, font_size=12)
add_box(slide, 7.35, 2.35, 1.9, 0.8, "WhatsApp", fill=RGBColor(240, 249, 255), line=BLUE, font_size=12)
add_box(slide, 9.55, 2.35, 1.7, 0.8, "电话", fill=RGBColor(240, 249, 255), line=BLUE, font_size=12)
add_box(slide, 11.45, 2.35, 1.15, 0.8, "表单", fill=RGBColor(240, 249, 255), line=BLUE, font_size=12)
add_arrow(slide, 2.9, 2.75, 3.18, 2.75, color=GREEN)
add_arrow(slide, 5.2, 2.75, 5.48, 2.75, color=GREEN)
add_arrow(slide, 7.1, 2.75, 7.32, 2.75, color=BLUE)
add_arrow(slide, 9.3, 2.75, 9.52, 2.75, color=BLUE)
add_arrow(slide, 11.2, 2.75, 11.42, 2.75, color=BLUE)

add_bullet_list(slide, 0.75, 3.55, 4.0, 2.7, "联系方式建议", [
    "微信作为主入口，适合家长快速沟通",
    "WhatsApp 作为备用通道，适合海外家长",
    "二维码用于扫码加好友，减少跳转成本",
], accent=GREEN)
add_bullet_list(slide, 4.95, 3.55, 4.0, 2.7, "页面呈现方式", [
    "首页底部做快捷按钮",
    "联系页放完整说明与二维码",
    "提交预约后提示优先通过微信 / WhatsApp 联系",
], accent=BLUE)
add_card(slide, 9.15, 3.55, 3.4, 2.7, "当前已填信息", "微信昵称：IvyEuro学校刘老师\n微信号：sophialiu0826\nWhatsApp：+34-625-627-022", fill=WHITE, accent=AMBER, title_size=14, body_size=12)

# Slide 7 items to refine
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_header(slide, "待补强内容", "哪些内容还需要真实材料")
add_title_block(slide, "后续优先级", "现在结构已齐，下一步重点是把模板句换成真实内容。")

add_card(slide, 0.75, 2.0, 3.0, 1.55, "优先级高", "老师真实照片、真实资历、擅长科目", fill=WHITE, accent=RED)
add_card(slide, 3.95, 2.0, 3.0, 1.55, "优先级高", "家长真实评价、案例、学习前后对比", fill=WHITE, accent=RED)
add_card(slide, 7.15, 2.0, 3.0, 1.55, "优先级中", "新闻动态、活动通知、开班消息", fill=WHITE, accent=AMBER)
add_card(slide, 10.35, 2.0, 2.25, 1.55, "优先级中", "课程成果表达", fill=WHITE, accent=AMBER)
add_bullet_list(slide, 0.75, 4.0, 5.95, 2.3, "建议替换的占位句", [
    "“这里可以展示……”",
    "“可用于展示……”",
    "“后续可继续补充……”",
], accent=AMBER)
add_bullet_list(slide, 6.9, 4.0, 5.7, 2.3, "建议保留的真实信息", [
    "机构名、地址、电话、营业时间",
    "课程大类与联系方式",
    "课程表页和联系页结构",
], accent=GREEN)

# Slide 8 summary
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, NAVY)
add_header(slide, "结论与下一步", "交付 / 补强 / 上线")

box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.2), Inches(11.8), Inches(5.4))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(17, 27, 44)
box.line.color.rgb = RGBColor(48, 62, 84)

tx = slide.shapes.add_textbox(Inches(1.1), Inches(1.6), Inches(10.8), Inches(4.6))
tf = tx.text_frame
tf.word_wrap = True
for i, line in enumerate([
    "1. 当前网站已完成结构重建，核心业务内容已迁移。",
    "2. 暑期课程已改为季节性公告，避免长期占用首页。",
    "3. 微信、WhatsApp、二维码、电话、表单都已经统一。",
    "4. 后续最重要的是补老师真实资料、家长真实评价和真实新闻。",
    "5. 完成后即可进入最终上线前检查与部署。",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run()
    r.text = line
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(17 if i == 0 else 15)
    r.font.bold = i == 0
    r.font.color.rgb = WHITE
    p.space_after = Pt(14)

foot = slide.shapes.add_textbox(Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.4))
fp = foot.text_frame.paragraphs[0]
rr = fp.add_run()
rr.text = "IvyEuro Education · 网站改造分析报告"
rr.font.name = "Microsoft YaHei"
rr.font.size = Pt(11)
rr.font.color.rgb = RGBColor(205, 214, 224)
fp.alignment = PP_ALIGN.RIGHT

prs.save(OUT)
print(OUT)

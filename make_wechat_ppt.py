from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "IvyEuro_Education_微信转发简版.pptx"

BLUE = RGBColor(27, 109, 255)
NAVY = RGBColor(11, 18, 32)
GREEN = RGBColor(34, 163, 92)
AMBER = RGBColor(245, 158, 11)
RED = RGBColor(239, 68, 68)
TEXT = RGBColor(16, 32, 51)
MUTED = RGBColor(90, 107, 127)
LIGHT = RGBColor(244, 247, 251)
WHITE = RGBColor(255, 255, 255)


def bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def header(slide, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.82))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.15), Inches(12), Inches(0.32))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = WHITE


def textbox(slide, x, y, w, h, text, size=14, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    p = t.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return t


def card(slide, x, y, w, h, title, body, accent=BLUE):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = WHITE
    shp.line.color.rgb = RGBColor(220, 227, 235)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.12), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    textbox(slide, x + 0.22, y + 0.15, w - 0.35, 0.28, title, size=15, bold=True)
    textbox(slide, x + 0.22, y + 0.52, w - 0.35, h - 0.58, body, size=11, color=MUTED)


def arrow(slide, x1, y1, x2, y2, color=RGBColor(145, 155, 170)):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(1.5)
    conn.line.end_arrowhead = True
    return conn


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Slide 1
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, NAVY)
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(0.8), Inches(2.7), Inches(0.48))
box.fill.solid(); box.fill.fore_color.rgb = BLUE; box.line.fill.background()
textbox(s, 0.75, 0.88, 2.7, 0.2, "微信转发简版", size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
textbox(s, 0.75, 1.35, 7.9, 1.7, "IvyEuro Education\n网站改造简报", size=28, color=WHITE, bold=True)
textbox(s, 0.75, 3.15, 7.5, 0.8, "从老站内容迁移到新官网，重点解决页面结构、咨询入口和内容更新问题。", size=14, color=RGBColor(204, 214, 224))
card(s, 8.75, 1.2, 3.8, 1.3, "项目状态", "已完成首页、多页面、新闻详情、联系通道和移动端适配。", accent=GREEN)
card(s, 8.75, 2.75, 3.8, 1.3, "转发用途", "适合微信里快速给客户、老板或家长看概览。", accent=AMBER)
card(s, 8.75, 4.3, 3.8, 1.3, "最新信息", f"生成日期：{date.today().isoformat()}", accent=BLUE)

# Slide 2
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header(s, "改造亮点")
textbox(s, 0.65, 1.1, 11.5, 0.55, "用最少的文字说明这次改造做了什么", size=26, bold=True)
card(s, 0.65, 1.95, 3.9, 1.55, "1. 页面结构更清晰", "首页 + 服务项目 + 关于我们 + 老师简介 + 课程表 + 家长好评 + 新闻动态 + 联系我们。", BLUE)
card(s, 4.72, 1.95, 3.9, 1.55, "2. 咨询入口更明确", "微信昵称、微信号、WhatsApp、二维码和表单都统一到联系区。", GREEN)
card(s, 8.79, 1.95, 3.9, 1.55, "3. 临时内容可自动隐藏", "暑期班只保留季节公告，不会长期压在首页。", AMBER)
card(s, 0.65, 3.85, 3.9, 1.55, "4. 新闻页更像正式栏目", "列表 + 详情页，方便后续持续发布通知和文章。", BLUE)
card(s, 4.72, 3.85, 3.9, 1.55, "5. 移动端能正常浏览", "导航、按钮、联系页都做了响应式适配。", GREEN)
card(s, 8.79, 3.85, 3.9, 1.55, "6. 视觉更像机构官网", "学校场景图、课程块和配色统一，正式感更强。", AMBER)

# Slide 3
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header(s, "内容取舍")
textbox(s, 0.65, 1.1, 11.6, 0.5, "哪些保留，哪些替换，哪些要补强", size=26, bold=True)
card(s, 0.65, 1.95, 3.75, 1.55, "保留", "机构名、地址、电话、营业时间、课程大类、微信 / WhatsApp。", GREEN)
card(s, 4.8, 1.95, 3.75, 1.55, "替换", "老师占位卡、家长好评模板句、新闻/文章草稿标题。", AMBER)
card(s, 8.95, 1.95, 3.75, 1.55, "补强", "真实老师资料、真实评价、真实新闻、真实学习成果。", BLUE)
card(s, 2.0, 4.1, 3.8, 1.25, "要删除", "过期的暑期首页大块展示。", RED)
card(s, 7.3, 4.1, 3.8, 1.25, "要季节化", "暑期课程改成公告，到期自动隐藏。", GREEN)
arrow(s, 2.5, 3.5, 2.5, 4.05, color=RED)
arrow(s, 9.2, 3.5, 9.2, 4.05, color=GREEN)

# Slide 4
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, NAVY)
header(s, "联系方式与下一步")
textbox(s, 0.65, 1.1, 11.7, 0.55, "微信优先，WhatsApp 备用", size=25, color=WHITE, bold=True)
card(s, 0.65, 1.95, 4.0, 1.2, "微信", "昵称：IvyEuro学校刘老师\n微信号：sophialiu0826", accent=GREEN)
card(s, 4.85, 1.95, 4.0, 1.2, "WhatsApp", "直接点击联系，适合海外家长。", accent=BLUE)
card(s, 9.05, 1.95, 3.55, 1.2, "二维码", "扫码后可继续咨询。", accent=AMBER)
textbox(s, 0.65, 3.55, 11.9, 1.0, "下一步建议：补真实老师、家长反馈、新闻内容，然后就可以做最终上线前检查。", size=16, color=WHITE, bold=True)
textbox(s, 0.65, 5.1, 11.9, 0.7, "IvyEuro Education · 微信转发简版", size=11, color=RGBColor(205, 214, 224), align=PP_ALIGN.RIGHT)

prs.save(OUT)
print(OUT)
